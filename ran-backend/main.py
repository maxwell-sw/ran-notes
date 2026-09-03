from fastapi import FastAPI, Request, HTTPException
from fastapi.responses import FileResponse, StreamingResponse
from fastapi import UploadFile, File, Form
from typing import List
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Literal, Optional
from openai import OpenAI
from dotenv import load_dotenv
import os
import re
import json
import base64
import asyncio
import subprocess
import shutil
import sqlite3
import uuid
import logging
from datetime import datetime
from collections import Counter
from tempfile import TemporaryDirectory
from pathlib import Path

try:
    from docx import Document
except ImportError:
    Document = None

load_dotenv()

MODEL_TIMEOUT_SECONDS = float(os.getenv("MODEL_TIMEOUT_SECONDS", "300"))
TRANSCRIPT_CHUNK_CHARS = int(os.getenv("TRANSCRIPT_CHUNK_CHARS", "6000"))
# Submit a small batch at a time: it materially speeds up long transcripts
# while avoiding a burst of every chunk against the provider at once.
TRANSCRIPT_MAX_CONCURRENCY = max(1, int(os.getenv("TRANSCRIPT_MAX_CONCURRENCY", "2")))
VISION_INPUT_ENABLED = os.getenv("VISION_INPUT_ENABLED", "true").strip().lower() not in {"0", "false", "no"}
VISION_MAX_IMAGES = max(0, int(os.getenv("VISION_MAX_IMAGES", "12")))

client = OpenAI(
    api_key=os.getenv("OPENAI_API_KEY"),
    base_url=os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1"),
    timeout=MODEL_TIMEOUT_SECONDS,
    max_retries=1,
)

# 模型选择：在 .env 中设置 OPENAI_MODEL 可覆盖默认值。推荐见 模型推荐.md
MODEL_NAME = os.getenv("OPENAI_MODEL", "gpt-4o")

app = FastAPI()
logger = logging.getLogger("ran_notes")

LIBRARY_DB = Path(__file__).with_name("ran_notes_library.sqlite3")
LIBRARY_FILES_ROOT = Path(__file__).with_name("ran_notes_assets")
TRIAL_ASSETS_ROOT = Path(__file__).with_name("trial_assets")

# 试用案例只暴露这里列出的文件，不把服务器上的任意路径提供给浏览器。
TRIAL_CASE = {
    "id": "pig-cycle-policy",
    "title": "中国猪周期运行机制与政策调控",
    "description": "含研究汇报 PPT、文献 PDF 与已转写的组会语音，适合首次体验完整流程。",
    "meeting_type": "人文社科",
    "discipline": "农业经济",
    "participants": ["李教授", "张同学", "王同学", "陈同学"],
    "reporters": ["张同学"],
    "materials": [
        {"id": "pig-ppt", "filename": "猪周期研究汇报PPT.pptx", "role": "material"},
        {"id": "pig-paper", "filename": "生猪价格指数保险能否抑制猪周期_廖朴.pdf", "role": "material"},
        {"id": "pig-transcript", "filename": "猪周期组会语音转写.docx", "role": "transcript"},
    ],
    # CAJ 是知网格式，当前解析链不支持；仍随试用包保留，供用户下载查看原件。
    "reference_download": {"id": "pig-caj", "filename": "生猪价格指数保险能否抑制猪周期_廖朴.caj"},
}


def _db():
    conn = sqlite3.connect(LIBRARY_DB)
    conn.row_factory = sqlite3.Row
    return conn


def _init_library():
    with _db() as conn:
        conn.executescript("""
        CREATE TABLE IF NOT EXISTS meetings (
            id TEXT PRIMARY KEY, meeting_time TEXT, meeting_type TEXT, discipline TEXT,
            summary TEXT, data_json TEXT NOT NULL, created_at TEXT NOT NULL
        );
        CREATE TABLE IF NOT EXISTS action_items (
            id TEXT PRIMARY KEY, meeting_id TEXT NOT NULL, reporter TEXT, description TEXT,
            owner TEXT, due TEXT, level TEXT, risk TEXT, status TEXT NOT NULL DEFAULT 'open',
            evidence_json TEXT, created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
            FOREIGN KEY(meeting_id) REFERENCES meetings(id)
        );
        """)
        # 兼容已创建过的资料库：补上行动项“是否已经回应”的闭环字段。
        action_columns = {row["name"] for row in conn.execute("PRAGMA table_info(action_items)")}
        if "response_status" not in action_columns:
            conn.execute("ALTER TABLE action_items ADD COLUMN response_status TEXT NOT NULL DEFAULT 'unanswered'")
        if "response_note" not in action_columns:
            conn.execute("ALTER TABLE action_items ADD COLUMN response_note TEXT NOT NULL DEFAULT ''")


@app.on_event("startup")
def init_library():
    _init_library()


def _flatten_actions(data: dict):
    actions = []
    for block in data.get("by_reporter") or []:
        reporter = str(block.get("reporter") or "待确认")
        for item in block.get("action_items") or []:
            if isinstance(item, dict):
                actions.append((reporter, item))
    for item in data.get("action_items") or []:
        if isinstance(item, dict):
            actions.append((str(item.get("owner") or "待确认"), item))
    return actions


def _safe_filename(filename: str) -> str:
    """Keep uploaded names inside the meeting archive directory."""
    return Path(str(filename or "attachment")).name or "attachment"


def _archive_preview_text(name: str, content: bytes) -> tuple[str, str]:
    """Return a short, local preview for files added directly in the library."""
    suffix = Path(name).suffix.lower()
    if suffix == ".pdf":
        return "pdf", extract_pdf_text_bytes(content)[:30000]
    if suffix in {".ppt", ".pptx"}:
        text, _ = extract_ppt_content_bytes(content)
        return "ppt", text[:30000]
    if suffix in {".doc", ".docx"}:
        return "word", extract_docx_text_bytes(content)[:30000]
    if suffix in {".txt", ".md"}:
        return "text", content.decode("utf-8", errors="replace")[:30000]
    return "file", "该文件可下载查看，暂不支持文本预览。"


def _meeting_row_with_data(meeting_id: str):
    with _db() as conn:
        row = conn.execute("SELECT * FROM meetings WHERE id=?", (meeting_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="未找到该场组会")
    try:
        data = json.loads(row["data_json"])
    except (json.JSONDecodeError, TypeError):
        data = {}
    return dict(row), data


def _save_meeting(data: dict, meeting_time: str, discipline: str, archive_inputs: Optional[List[dict]] = None):
    meeting_id = str(uuid.uuid4())
    now = datetime.now().isoformat(timespec="seconds")
    archive_dir = LIBRARY_FILES_ROOT / meeting_id
    archived_assets = []
    for source in archive_inputs or []:
        content = source.get("content") if isinstance(source, dict) else None
        if not isinstance(content, bytes):
            continue
        asset_id = str(uuid.uuid4())
        original_name = _safe_filename(source.get("name"))
        stored_name = f"{asset_id}_{original_name}"
        archive_dir.mkdir(parents=True, exist_ok=True)
        (archive_dir / stored_name).write_bytes(content)
        archived_assets.append({
            "id": asset_id,
            "name": original_name,
            "stored_name": stored_name,
            "type": str(source.get("type") or "file"),
            "owners": [str(owner).strip() for owner in (source.get("owners") or []) if str(owner).strip()] or ["其他材料"],
            "preview_text": str(source.get("preview_text") or "")[:30000],
        })
    # File metadata lives with the meeting record; the actual bytes live outside SQLite.
    data["archived_assets"] = archived_assets
    with _db() as conn:
        conn.execute(
            "INSERT INTO meetings VALUES (?, ?, ?, ?, ?, ?, ?)",
            (meeting_id, meeting_time, data.get("meeting_type", ""), discipline, data.get("summary", ""), json.dumps(data, ensure_ascii=False), now),
        )
        for reporter, item in _flatten_actions(data):
            action_id = str(uuid.uuid4())
            conn.execute(
                """INSERT INTO action_items
                (id, meeting_id, reporter, description, owner, due, level, risk, status,
                 evidence_json, response_status, response_note, created_at, updated_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, 'open', ?, 'unanswered', '', ?, ?)""",
                (action_id, meeting_id, reporter, item.get("description", ""), item.get("owner", reporter), item.get("due", "未明确"), item.get("level", "YELLOW"), item.get("risk", ""), json.dumps(item.get("evidence") or [], ensure_ascii=False), now, now),
            )
    return meeting_id


def _focus_terms(data: dict):
    terms = []
    for block in data.get("by_reporter") or []:
        for item in block.get("advisor_feedback") or []:
            terms.extend(item.get("focus_tags") or [])
    terms.extend(data.get("mentor_focus") or [])
    return [str(term).strip() for term in terms if str(term).strip()]

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class MeetingInfo(BaseModel):
    time: str
    topic: str
    roles: List[str]


class GenerateRecordRequest(BaseModel):
    meetingInfo: MeetingInfo
    transcript: str
    pptSummary: Optional[str] = ""
    papersSummary: Optional[str] = ""


@app.post("/generate-record")
def extract_pdf_text(file: UploadFile) -> str:
    try:
        reader = PdfReader(BytesIO(file.file.read()))
        texts = []
        for i, page in enumerate(reader.pages):
            txt = page.extract_text() or ""
            if txt.strip():
                texts.append(f"[PDF 第{i+1}页]\n{txt}")
        return "\n\n".join(texts)
    except Exception as e:
        return f"[PDF 解析失败: {str(e)}]"

def extract_ppt_text(file: UploadFile) -> str:
    try:
        prs = Presentation(BytesIO(file.file.read()))
        texts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                texts.append(f"[PPT 第{i+1}页]\n" + "\n".join(slide_text))
        return "\n\n".join(texts)
    except Exception as e:
        return f"[PPT 解析失败: {str(e)}]"


def extract_docx_text(file: UploadFile) -> str:
    if Document is None:
        return "[Word 解析失败: 未安装 python-docx]"
    try:
        file_bytes = BytesIO(file.file.read())
        doc = Document(file_bytes)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[Word 解析失败: {str(e)}]"


def extract_pdf_text_bytes(content: bytes) -> str:
    try:
        reader = PdfReader(BytesIO(content))
        texts = []
        for i, page in enumerate(reader.pages):
            txt = page.extract_text() or ""
            if txt.strip():
                texts.append(f"[PDF 第{i+1}页]\n{txt}")
        return "\n\n".join(texts)
    except Exception as e:
        return f"[PDF 解析失败: {str(e)}]"


def extract_pdf_pages_with_images(content: bytes) -> list:
    """按页提取 PDF：每页返回 { page, text, image }，image 为 base64 数据 URL。无 PyMuPDF 时仅返回 text。"""
    try:
        import fitz  # type: ignore  # PyMuPDF
    except ImportError:
        try:
            reader = PdfReader(BytesIO(content))
            return [
                {"page": i + 1, "text": (p.extract_text() or "").strip(), "image": None}
                for i, p in enumerate(reader.pages)
            ]
        except Exception:
            return []
    try:
        doc = fitz.open(stream=content, filetype="pdf")  # type: ignore
        pages = []
        for i in range(len(doc)):
            page = doc[i]
            text = (page.get_text() or "").strip()
            pix = page.get_pixmap(dpi=144)
            png_bytes = pix.tobytes("png")
            b64 = base64.b64encode(png_bytes).decode("ascii")
            pages.append({
                "page": i + 1,
                "text": text,
                "image": f"data:image/png;base64,{b64}",
            })
        doc.close()
        return pages
    except Exception:
        try:
            reader = PdfReader(BytesIO(content))
            return [
                {"page": i + 1, "text": (p.extract_text() or "").strip(), "image": None}
                for i, p in enumerate(reader.pages)
            ]
        except Exception:
            return []


def extract_docx_text_bytes(content: bytes) -> str:
    if Document is None:
        return "[Word 解析失败: 未安装 python-docx]"
    try:
        doc = Document(BytesIO(content))
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[Word 解析失败: {str(e)}]"


def extract_ppt_content_bytes(content: bytes):
    try:
        prs = Presentation(BytesIO(content))
        texts = []
        slides_meta = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            images = []
            text_boxes = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                    text_boxes.append({
                        "text": shape.text,
                        "left": int(shape.left), "top": int(shape.top),
                        "width": int(shape.width), "height": int(shape.height),
                    })
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    ext = (img.ext or "png").lower()
                    mime = "image/png"
                    if ext in ["jpg", "jpeg"]:
                        mime = "image/jpeg"
                    b64 = base64.b64encode(img.blob).decode("ascii")
                    images.append(f"data:{mime};base64,{b64}")
            if slide_text:
                texts.append(f"[PPT 第{i+1}页]\n" + "\n".join(slide_text))
            slides_meta.append(
                {
                    "page": i + 1,
                    "images": images,
                    "text_boxes": text_boxes,
                    "slide_width": int(prs.slide_width),
                    "slide_height": int(prs.slide_height),
                }
            )
        return "\n\n".join(texts), slides_meta
    except Exception as e:
        return f"[PPT 解析失败: {str(e)}]", []


def _find_soffice():
    """查找本机 soffice 可执行文件，优先 PATH，再常见路径（Windows / macOS / Linux）。"""
    exe = shutil.which("soffice") or shutil.which("libreoffice")
    if exe:
        return exe

    candidates = []

    # Windows 常见安装路径
    import sys
    if sys.platform == "win32":
        win_candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        # 兼容中文/英文系统的 Program Files 环境变量
        for env_var in ("PROGRAMFILES", "PROGRAMFILES(X86)", "PROGRAMW6432"):
            pf = os.environ.get(env_var)
            if pf:
                win_candidates.append(os.path.join(pf, "LibreOffice", "program", "soffice.exe"))
        candidates.extend(win_candidates)

    # macOS 常见路径
    candidates += [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/local/bin/soffice",
        "/opt/homebrew/bin/soffice",
    ]
    # macOS 模糊匹配（版本号不固定时）
    applications = Path("/Applications")
    if applications.exists():
        for app in applications.glob("LibreOffice*.app"):
            p = app / "Contents/MacOS/soffice"
            if p.exists():
                candidates.insert(0, str(p))

    for p in candidates:
        if os.path.isfile(p) and os.access(p, os.X_OK):
            return p
    return None


def _run_libreoffice_convert(ppt_path: Path, outdir: Path) -> bool:
    """使用 LibreOffice 将 PPT 转为 PDF，成功返回 True。"""
    ppt_abs = ppt_path.resolve()
    outdir_abs = outdir.resolve()
    base = ["--headless", "--convert-to", "pdf", str(ppt_abs), "--outdir", str(outdir_abs)]
    env = os.environ.copy()
    env.setdefault("HOME", str(outdir_abs))
    env["SAL_USE_VCLPLUGIN"] = "gen"
    candidates = []
    soffice_exe = _find_soffice()
    if soffice_exe:
        candidates.append([soffice_exe] + base)
    candidates.extend([
        ["libreoffice"] + base,
        ["soffice"] + base,
    ])
    for cmd in candidates:
        try:
            subprocess.run(
                cmd,
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=120,
                env=env,
                cwd=str(outdir_abs),
            )
            return True
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError):
            continue
    return False


def _render_ppt_slide_previews(content: bytes) -> list:
    """快速绘制每一页的可视化预览，不依赖 LibreOffice，也不会阻塞会议生成。"""
    try:
        from PIL import Image, ImageDraw, ImageFont  # type: ignore
        prs = Presentation(BytesIO(content))
        width = 1600
        height = max(1, round(width * prs.slide_height / prs.slide_width))
        x_scale, y_scale = width / prs.slide_width, height / prs.slide_height
        font_paths = [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Light.ttc",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        ]

        def font(size):
            for path in font_paths:
                try:
                    return ImageFont.truetype(path, size)
                except OSError:
                    continue
            return ImageFont.load_default()

        previews = []
        for slide in prs.slides:
            canvas = Image.new("RGB", (width, height), "white")
            draw = ImageDraw.Draw(canvas)
            for shape in slide.shapes:
                left = int(shape.left * x_scale)
                top = int(shape.top * y_scale)
                right = int((shape.left + shape.width) * x_scale)
                bottom = int((shape.top + shape.height) * y_scale)
                if right <= left or bottom <= top:
                    continue
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        picture = Image.open(BytesIO(shape.image.blob)).convert("RGB")
                        picture.thumbnail((right - left, bottom - top))
                        canvas.paste(picture, (left + max(0, (right - left - picture.width) // 2), top + max(0, (bottom - top - picture.height) // 2)))
                        continue
                except Exception:
                    pass
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    draw.rectangle((left, top, right, bottom), outline="#cbd5e1", width=1)
                    size = max(14, min(34, int((bottom - top) / max(2, text.count("\n") + 2))))
                    draw.multiline_text((left + 8, top + 6), text, fill="#0f172a", font=font(size), spacing=5)
            buffer = BytesIO()
            canvas.save(buffer, format="PNG", optimize=True)
            previews.append("data:image/png;base64," + base64.b64encode(buffer.getvalue()).decode("ascii"))
        return previews
    except Exception:
        return []


def _render_ppt_with_libreoffice(content: bytes):
    try:
        import fitz  # type: ignore
    except ImportError:
        return []

    # Windows 下 soffice/fitz 偶发短暂占用临时 PDF；Python 3.10+ 可用 ignore_cleanup_errors，3.9 用普通 TemporaryDirectory
    with TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)
        ppt_path = tmpdir_path / "slides.pptx"
        pdf_path = tmpdir_path / "slides.pdf"
        ppt_path.write_bytes(content)

        if not _run_libreoffice_convert(ppt_path, tmpdir_path):
            return []

        if not pdf_path.exists():
            candidates = list(tmpdir_path.glob("*.pdf"))
            if not candidates:
                return []
            pdf_path = candidates[0]

        doc = None
        try:
            doc = fitz.open(pdf_path)  # type: ignore
        except Exception:
            return []

        images = []
        try:
            for page_index in range(len(doc)):
                page = doc[page_index]
                pix = page.get_pixmap(dpi=144)
                png_bytes = pix.tobytes("png")
                b64 = base64.b64encode(png_bytes).decode("ascii")
                images.append(f"data:image/png;base64,{b64}")
        finally:
            # Windows 下若不及时 close，临时目录清理会出现 WinError 32（文件占用）
            try:
                if doc is not None:
                    doc.close()
            except Exception:
                pass

        return images


def render_ppt_to_page_images(content: bytes):
    # 先用快速页级预览保证交互即时、每页可见；Office 转换仅作为没有可用预览时的兜底。
    previews = _render_ppt_slide_previews(content)
    return previews if previews else _render_ppt_with_libreoffice(content)


def _parse_json_list(value: str) -> list:
    """解析表单中的 JSON 数组；不合法时返回空列表。"""
    try:
        parsed = json.loads(value or "[]")
        return parsed if isinstance(parsed, list) else []
    except (TypeError, json.JSONDecodeError):
        return []


@app.post("/prepare-transcript")
async def prepare_transcript(request: Request):
    """把原始 ASR 转写整理为可核对的会议转写，不直接生成纪要。"""
    form = await request.form()
    raw_transcript = (form.get("raw_transcript") or "").strip()
    speaker_instruction = (form.get("speaker_instruction") or "").strip()
    transcript_file = form.get("transcript_file")

    if getattr(transcript_file, "filename", None):
        filename = transcript_file.filename.lower()
        content = await transcript_file.read()
        if filename.endswith(".txt"):
            file_text = content.decode("utf-8", errors="replace")
        elif filename.endswith(".docx"):
            file_text = extract_docx_text_bytes(content)
        else:
            raise HTTPException(status_code=400, detail="转写稿仅支持 TXT 或 DOCX 格式。")
        raw_transcript = (raw_transcript + "\n\n" + file_text).strip()

    if not raw_transcript:
        raise HTTPException(status_code=400, detail="请提供原始语音转写稿。")

    system_prompt = """你是科研组会的语音转写编辑。你的任务仅是把 ASR 原始稿整理成可供人工核对的文本，不要生成会议纪要。

必须遵守：
1. 不得编造、补全或改写未说出的研究事实、数据、结论、人名、文献或任务；听不清、语义残缺处保留原意并可标为「[听辨不清]」。
2. 按发言轮次断句，统一明显重复的标签格式（如 SPEAKER_00、发言人1），修正常见标点、断句和无意义口头填充，但不能删除有信息价值的内容。
3. 仅将用户明确说明的标签映射为姓名/角色。未被明确说明的身份必须保留为原始标签或中性标签，绝不可猜测其姓名、导师身份或汇报人身份。
4. 输出中每段都应有「姓名或标签：内容」的前缀，方便后续人工修改。
5. 同时提取可供用户核对的会议信息：日期和时间仅在原文有完整、明确的公历日期与时间时填写 YYYY-MM-DD HH:MM；相对日期、缺少年份或不明确时留空。会议类型只能是「人文社科、实验工科、理论数理、通用组会」之一或空。学科、参与人和汇报人只能依据原文和用户明确说明；汇报人仅列出实际做研究汇报的人，不确定时留空。
6. 只返回 JSON：{"cleaned_transcript":"...","speakers":[{"label":"原始标签","name":"整理后显示名称","confidence":"provided|unconfirmed"}],"meeting_suggestions":{"time":"YYYY-MM-DD HH:MM 或空","meeting_type":"四种类型之一或空","discipline":"...或空","participants":["..."],"reporters":["..."]},"notes":["..."]}。speakers 只列实际在转写中出现的标签；无明确映射时 name 等于 label，confidence 为 unconfirmed。"""
    try:
        response = await asyncio.to_thread(
            client.chat.completions.create,
            model=MODEL_NAME,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": json.dumps({"speaker_instruction": speaker_instruction or "（未提供；不要推测身份）", "raw_transcript": raw_transcript}, ensure_ascii=False)},
            ],
            response_format={"type": "json_object"},
        )
        data = json.loads(response.choices[0].message.content)
    except Exception as exc:
        provider_status = getattr(exc, "status_code", None)
        logger.warning("prepare-transcript model request failed: type=%s status=%s message=%s", type(exc).__name__, provider_status, str(exc)[:300])
        if provider_status == 401:
            detail = "模型服务鉴权失败，请检查后端的 OPENAI_API_KEY。"
        elif provider_status == 403:
            detail = f"模型“{MODEL_NAME}”当前不可用，请检查模型名称或服务权限。"
        elif provider_status == 429:
            detail = "模型服务请求过于频繁或额度不足，请稍后再试。"
        elif provider_status in {500, 502, 503, 504}:
            detail = "模型服务暂时没有成功回应（上游服务错误），请稍后重试。"
        elif "timeout" in str(exc).lower() or "timed out" in str(exc).lower():
            detail = "模型整理超时，请稍后重试或缩短单次转写稿。"
        else:
            detail = "转写整理未完成，请检查网络和模型服务状态后重试。"
        raise HTTPException(status_code=502, detail=detail) from exc

    cleaned_transcript = str(data.get("cleaned_transcript") or "").strip() or raw_transcript
    speakers = []
    for item in data.get("speakers") or []:
        if not isinstance(item, dict):
            continue
        label = str(item.get("label") or "").strip()
        name = str(item.get("name") or label).strip()
        if label or name:
            speakers.append({"label": label or name, "name": name or label, "confidence": "provided" if item.get("confidence") == "provided" else "unconfirmed"})
    suggestion_data = data.get("meeting_suggestions") if isinstance(data.get("meeting_suggestions"), dict) else {}
    suggested_time = str(suggestion_data.get("time") or "").strip()
    if not re.fullmatch(r"\d{4}-\d{2}-\d{2}\s+\d{2}:\d{2}", suggested_time):
        suggested_time = ""
    suggested_type = str(suggestion_data.get("meeting_type") or "").strip()
    if suggested_type not in {"人文社科", "实验工科", "理论数理", "通用组会"}:
        suggested_type = ""
    def _suggested_names(value):
        return [str(name).strip() for name in value if str(name).strip()][:20] if isinstance(value, list) else []
    meeting_suggestions = {
        "time": suggested_time,
        "meeting_type": suggested_type,
        "discipline": str(suggestion_data.get("discipline") or "").strip()[:80],
        "participants": _suggested_names(suggestion_data.get("participants")),
        "reporters": _suggested_names(suggestion_data.get("reporters")),
    }
    return {"original_transcript": raw_transcript, "cleaned_transcript": cleaned_transcript, "speakers": speakers, "meeting_suggestions": meeting_suggestions, "notes": [str(note) for note in (data.get("notes") or []) if str(note).strip()][:5]}


def _split_transcript_for_model(text: str, limit: int = TRANSCRIPT_CHUNK_CHARS) -> list:
    """Split on ASR turns where possible so a completed segment is a real unit of work."""
    if len(text) <= limit:
        return [text]
    chunks, current, current_size = [], [], 0
    for line in text.splitlines(keepends=True):
        if current and current_size + len(line) > limit:
            chunks.append("".join(current).strip())
            current, current_size = [], 0
        # A single pathological ASR line should not create a huge model request.
        while len(line) > limit:
            if current:
                chunks.append("".join(current).strip())
                current, current_size = [], 0
            chunks.append(line[:limit].strip())
            line = line[limit:]
        current.append(line)
        current_size += len(line)
    if current:
        chunks.append("".join(current).strip())
    return [chunk for chunk in chunks if chunk]


def _model_error_detail(exc: Exception) -> str:
    status = getattr(exc, "status_code", None)
    if status == 401:
        return "模型服务鉴权失败，请检查后端的 OPENAI_API_KEY。"
    if status == 403:
        return f"模型“{MODEL_NAME}”当前不可用，请检查模型名称或服务权限。"
    if status == 429:
        return "当前模型的上游服务暂时限流（429），不是转写格式问题。请等待约 30 秒后重试，或在 .env 更换可用模型。"
    if status in {500, 502, 503, 504}:
        return "模型服务暂时没有成功回应（上游服务错误），请稍后重试。"
    if "timeout" in str(exc).lower() or "timed out" in str(exc).lower():
        return "模型整理超时，请稍后重试或缩短单次转写稿。"
    return "转写整理未完成，请检查网络和模型服务状态后重试。"


def _normalise_chunk_result(data: dict, fallback_text: str) -> dict:
    speakers = []
    for item in data.get("speakers") or []:
        if not isinstance(item, dict):
            continue
        label = str(item.get("label") or "").strip()
        name = str(item.get("name") or label).strip()
        if label or name:
            speakers.append({"label": label or name, "name": name or label, "confidence": "provided" if item.get("confidence") == "provided" else "unconfirmed"})
    suggestions = data.get("meeting_suggestions") if isinstance(data.get("meeting_suggestions"), dict) else {}
    suggested_time = str(suggestions.get("time") or "").strip()
    if not re.fullmatch(r"\d{4}-\d{2}-\d{2}\s+\d{2}:\d{2}", suggested_time):
        suggested_time = ""
    suggested_type = str(suggestions.get("meeting_type") or "").strip()
    if suggested_type not in {"人文社科", "实验工科", "理论数理", "通用组会"}:
        suggested_type = ""
    def names(value):
        return [str(name).strip() for name in value if str(name).strip()][:20] if isinstance(value, list) else []
    return {
        "cleaned_transcript": str(data.get("cleaned_transcript") or "").strip() or fallback_text,
        "speakers": speakers,
        "meeting_suggestions": {
            "time": suggested_time,
            "meeting_type": suggested_type,
            "discipline": str(suggestions.get("discipline") or "").strip()[:80],
            "participants": names(suggestions.get("participants")),
            "reporters": names(suggestions.get("reporters")),
        },
        "notes": [str(note).strip() for note in (data.get("notes") or []) if str(note).strip()][:5],
    }


def _merge_chunk_results(raw_transcript: str, results: list) -> dict:
    ordered = [item for item in sorted(results, key=lambda value: value[0])]
    speakers, seen_speakers, notes = [], set(), []
    suggestions = {"time": "", "meeting_type": "", "discipline": "", "participants": [], "reporters": []}
    for _, result in ordered:
        for speaker in result.get("speakers") or []:
            key = f"{speaker.get('label')}|{speaker.get('name')}"
            if key not in seen_speakers:
                seen_speakers.add(key)
                speakers.append(speaker)
        current = result.get("meeting_suggestions") or {}
        for field in ("time", "meeting_type", "discipline"):
            if not suggestions[field] and current.get(field):
                suggestions[field] = current[field]
        for field in ("participants", "reporters"):
            for name in current.get(field) or []:
                if name not in suggestions[field]:
                    suggestions[field].append(name)
        for note in result.get("notes") or []:
            if note not in notes and len(notes) < 5:
                notes.append(note)
    return {
        "original_transcript": raw_transcript,
        "cleaned_transcript": "\n\n".join(item[1]["cleaned_transcript"] for item in ordered),
        "speakers": speakers,
        "meeting_suggestions": suggestions,
        "notes": notes,
    }


@app.post("/prepare-transcript-stream")
async def prepare_transcript_stream(request: Request):
    """Stream actual per-segment completion events while the model edits an ASR transcript."""
    form = await request.form()
    raw_transcript = (form.get("raw_transcript") or "").strip()
    speaker_instruction = (form.get("speaker_instruction") or "").strip()
    transcript_file = form.get("transcript_file")
    if getattr(transcript_file, "filename", None):
        filename = transcript_file.filename.lower()
        content = await transcript_file.read()
        if filename.endswith(".txt"):
            file_text = content.decode("utf-8", errors="replace")
        elif filename.endswith(".docx"):
            file_text = extract_docx_text_bytes(content)
        else:
            raise HTTPException(status_code=400, detail="转写稿仅支持 TXT 或 DOCX 格式。")
        raw_transcript = (raw_transcript + "\n\n" + file_text).strip()
    if not raw_transcript:
        raise HTTPException(status_code=400, detail="请提供原始语音转写稿。")

    chunks = _split_transcript_for_model(raw_transcript)
    system_prompt = """你是科研组会的语音转写编辑。你正在处理完整转写的一部分，请只整理本段，不要补写本段之外的内容或会议纪要。
1. 不编造研究事实、数据、人名、文献或任务；听不清处保留原意并可标为「[听辨不清]」。
2. 按发言轮次断句，修正标点、明显重复标签和无意义口头填充；每段保留「姓名或标签：内容」前缀。
3. 仅使用用户明确给出的标签—身份映射；不确定身份保留原始标签。
4. 返回且仅返回 JSON：{"cleaned_transcript":"...","speakers":[{"label":"原始标签","name":"整理后显示名称","confidence":"provided|unconfirmed"}],"meeting_suggestions":{"time":"YYYY-MM-DD HH:MM 或空","meeting_type":"人文社科|实验工科|理论数理|通用组会或空","discipline":"...或空","participants":["..."],"reporters":["..."]},"notes":["..."]}。会议信息仅在本段有明确证据时填写。"""

    async def process_chunk(index: int, text: str):
        try:
            response = await asyncio.to_thread(
                client.chat.completions.create,
                model=MODEL_NAME,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": json.dumps({"speaker_instruction": speaker_instruction or "（未提供；不要推测身份）", "raw_transcript_segment": text, "segment_number": index + 1, "segment_total": len(chunks)}, ensure_ascii=False)},
                ],
                response_format={"type": "json_object"},
            )
            return index, _normalise_chunk_result(json.loads(response.choices[0].message.content), text), None
        except Exception as exc:
            logger.warning("prepare-transcript segment failed: index=%s type=%s status=%s message=%s", index, type(exc).__name__, getattr(exc, "status_code", None), str(exc)[:300])
            return index, None, _model_error_detail(exc)

    def sse(payload: dict) -> str:
        return f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"

    async def event_stream():
        total = len(chunks)
        yield sse({"event": "started", "total": total, "chunk_chars": TRANSCRIPT_CHUNK_CHARS, "concurrency": TRANSCRIPT_MAX_CONCURRENCY})
        results, failures, completed = [], [], 0
        # Work in small batches. If an upstream error occurs, cancel only the
        # other request in the current batch and never submit later batches.
        for batch_start in range(0, total, TRANSCRIPT_MAX_CONCURRENCY):
            batch = [
                asyncio.create_task(process_chunk(index, text))
                for index, text in enumerate(chunks[batch_start:batch_start + TRANSCRIPT_MAX_CONCURRENCY], start=batch_start)
            ]
            for task in asyncio.as_completed(batch):
                index, result, failure = await task
                if failure:
                    failures.append({"index": index + 1, "detail": failure})
                    for pending in batch:
                        if not pending.done():
                            pending.cancel()
                    await asyncio.gather(*batch, return_exceptions=True)
                    yield sse({"event": "segment_failed", "index": index + 1, "completed": completed, "total": total, "detail": failure})
                    yield sse({"event": "error", "detail": f"第{index + 1}段未提交成功：{failure} 已停止后续分段，避免继续触发上游限流。"})
                    return
                results.append((index, result))
                completed += 1
                yield sse({"event": "segment_complete", "index": index + 1, "completed": completed, "total": total})
        yield sse({"event": "result", "data": _merge_chunk_results(raw_transcript, results)})

    return StreamingResponse(event_stream(), media_type="text/event-stream", headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


class MeetingCreate(BaseModel):
    meeting_time: str
    meeting_type: str = ""
    discipline: str = ""
    summary: str = ""


class MeetingUpdate(BaseModel):
    meeting_time: Optional[str] = None
    meeting_type: Optional[str] = None
    discipline: Optional[str] = None
    summary: Optional[str] = None


@app.post("/library/meetings")
def create_library_meeting(meeting: MeetingCreate):
    meeting_id = str(uuid.uuid4())
    now = datetime.now().isoformat(timespec="seconds")
    data = {
        "meeting_type": meeting.meeting_type.strip(),
        "summary": meeting.summary.strip(),
        "basic_info": {"time": meeting.meeting_time.strip(), "discipline": meeting.discipline.strip()},
        "archived_assets": [],
    }
    with _db() as conn:
        conn.execute(
            "INSERT INTO meetings VALUES (?, ?, ?, ?, ?, ?, ?)",
            (meeting_id, meeting.meeting_time.strip(), meeting.meeting_type.strip(), meeting.discipline.strip(), meeting.summary.strip(), json.dumps(data, ensure_ascii=False), now),
        )
    return {"id": meeting_id}


@app.patch("/library/meetings/{meeting_id}")
def update_library_meeting(meeting_id: str, update: MeetingUpdate):
    row, data = _meeting_row_with_data(meeting_id)
    changes = {
        "meeting_time": row["meeting_time"], "meeting_type": row["meeting_type"],
        "discipline": row["discipline"], "summary": row["summary"],
    }
    for field in changes:
        value = getattr(update, field)
        if value is not None:
            changes[field] = value.strip()[:2000]
    data["meeting_type"] = changes["meeting_type"]
    data["summary"] = changes["summary"]
    basic = data.setdefault("basic_info", {})
    if isinstance(basic, dict):
        basic["time"] = changes["meeting_time"]
        basic["discipline"] = changes["discipline"]
    with _db() as conn:
        conn.execute(
            "UPDATE meetings SET meeting_time=?, meeting_type=?, discipline=?, summary=?, data_json=? WHERE id=?",
            (changes["meeting_time"], changes["meeting_type"], changes["discipline"], changes["summary"], json.dumps(data, ensure_ascii=False), meeting_id),
        )
    return {"id": meeting_id, **changes}


@app.delete("/library/meetings/{meeting_id}")
def delete_library_meeting(meeting_id: str):
    # Resolve the exact database row first; never derive a delete target from client path input.
    _meeting_row_with_data(meeting_id)
    with _db() as conn:
        conn.execute("DELETE FROM action_items WHERE meeting_id=?", (meeting_id,))
        conn.execute("DELETE FROM meetings WHERE id=?", (meeting_id,))
    archive_dir = (LIBRARY_FILES_ROOT / meeting_id).resolve()
    if archive_dir.parent == LIBRARY_FILES_ROOT.resolve() and archive_dir.is_dir():
        shutil.rmtree(archive_dir)
    return {"id": meeting_id, "deleted": True}


@app.post("/library/assets/{meeting_id}")
async def add_library_asset(meeting_id: str, file: UploadFile = File(...), owners: str = Form("")):
    _, data = _meeting_row_with_data(meeting_id)
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="不能添加空文件")
    if len(content) > 50 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="单个归档文件不能超过 50MB")
    original_name = _safe_filename(file.filename)
    asset_id = str(uuid.uuid4())
    stored_name = f"{asset_id}_{original_name}"
    asset_type, preview_text = _archive_preview_text(original_name, content)
    owner_list = [item.strip() for item in re.split(r"[,，、]", owners) if item.strip()] or ["其他材料"]
    archive_dir = LIBRARY_FILES_ROOT / meeting_id
    archive_dir.mkdir(parents=True, exist_ok=True)
    (archive_dir / stored_name).write_bytes(content)
    asset = {"id": asset_id, "name": original_name, "stored_name": stored_name, "type": asset_type, "owners": owner_list, "preview_text": preview_text}
    assets = data.setdefault("archived_assets", [])
    if not isinstance(assets, list):
        assets = data["archived_assets"] = []
    assets.append(asset)
    with _db() as conn:
        conn.execute("UPDATE meetings SET data_json=? WHERE id=?", (json.dumps(data, ensure_ascii=False), meeting_id))
    return {"id": asset_id, "name": original_name}


@app.delete("/library/assets/{meeting_id}/{asset_id}")
def delete_library_asset(meeting_id: str, asset_id: str):
    _, data = _meeting_row_with_data(meeting_id)
    assets = data.get("archived_assets") or []
    asset = next((item for item in assets if isinstance(item, dict) and item.get("id") == asset_id), None)
    if not asset:
        raise HTTPException(status_code=404, detail="未找到该材料")
    kept = [item for item in assets if not (isinstance(item, dict) and item.get("id") == asset_id)]
    data["archived_assets"] = kept
    with _db() as conn:
        conn.execute("UPDATE meetings SET data_json=? WHERE id=?", (json.dumps(data, ensure_ascii=False), meeting_id))
    path = (LIBRARY_FILES_ROOT / meeting_id / _safe_filename(asset.get("stored_name"))).resolve()
    root = (LIBRARY_FILES_ROOT / meeting_id).resolve()
    if path.parent == root and path.is_file():
        path.unlink()
    return {"id": asset_id, "deleted": True}


@app.get("/library/overview")
def library_overview():
    with _db() as conn:
        meetings = [dict(row) for row in conn.execute("SELECT id, meeting_time, meeting_type, discipline, summary, data_json FROM meetings ORDER BY meeting_time DESC, created_at DESC")]
        actions = [dict(row) for row in conn.execute("SELECT * FROM action_items ORDER BY created_at DESC")]
    focus = Counter()
    heat = Counter()
    trajectories = {}
    innovation_transfers = []
    mentor_briefs = []
    for meeting in meetings:
        data = json.loads(meeting.pop("data_json"))
        meeting["assets"] = [{
            "id": asset.get("id"),
            "name": asset.get("name"),
            "type": asset.get("type", "file"),
            "owners": asset.get("owners") or ["其他材料"],
        } for asset in data.get("archived_assets") or [] if isinstance(asset, dict)]
        for term in _focus_terms(data):
            focus[term] += 1
        day = (meeting.get("meeting_time") or "")[:10]
        if re.fullmatch(r"\d{4}-\d{2}-\d{2}", day):
            heat[day] += 1
        if data.get("mentor_brief"):
            mentor_briefs.append({"meeting_time": meeting.get("meeting_time"), "text": str(data["mentor_brief"])})
        for transfer in data.get("innovation_transfers") or []:
            if isinstance(transfer, dict):
                innovation_transfers.append({**transfer, "meeting_time": meeting.get("meeting_time")})
        for progress in data.get("student_progress") or []:
            if not isinstance(progress, dict):
                continue
            student = str(progress.get("student") or "").strip()
            if not student:
                continue
            trajectories.setdefault(student, []).append({
                "meeting_time": meeting.get("meeting_time"),
                "progress": str(progress.get("progress") or ""),
                "next_focus": str(progress.get("next_focus") or ""),
            })
    for action in actions:
        # 已完成或明确写入“已回应”的事项不会再显示为待回应；旧数据保持兼容。
        response = action.get("response_status") or "unanswered"
        action["response_status"] = response
        action["response_note"] = action.get("response_note") or ""
        action["needs_response"] = action.get("status") != "done" and response != "responded"
    return {
        "meetings": meetings,
        "heatmap": [{"date": day, "count": count} for day, count in heat.items()],
        "mentor_focus": [{"term": term, "count": count} for term, count in focus.most_common(12)],
        "actions": actions,
        "unanswered_actions": [item for item in actions if item["needs_response"]],
        "student_trajectories": [
            {"student": student, "updates": updates[:6]}
            for student, updates in sorted(trajectories.items())
        ],
        "innovation_transfers": innovation_transfers[:12],
        "mentor_briefs": mentor_briefs[:8],
    }


def _find_archived_asset(meeting_id: str, asset_id: str):
    with _db() as conn:
        row = conn.execute("SELECT data_json FROM meetings WHERE id=?", (meeting_id,)).fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="未找到该场组会")
    try:
        assets = json.loads(row["data_json"]).get("archived_assets") or []
    except (json.JSONDecodeError, TypeError):
        assets = []
    asset = next((item for item in assets if isinstance(item, dict) and item.get("id") == asset_id), None)
    if not asset:
        raise HTTPException(status_code=404, detail="未找到该材料")
    path = LIBRARY_FILES_ROOT / meeting_id / _safe_filename(asset.get("stored_name"))
    if not path.is_file():
        raise HTTPException(status_code=404, detail="归档文件不存在，可能已被移动")
    return asset, path


@app.get("/library/assets/{meeting_id}/{asset_id}/preview")
def preview_library_asset(meeting_id: str, asset_id: str):
    asset, _ = _find_archived_asset(meeting_id, asset_id)
    return {
        "name": asset.get("name"),
        "type": asset.get("type", "file"),
        "owners": asset.get("owners") or ["其他材料"],
        "text": asset.get("preview_text") or "该文件暂不支持文本预览，可打开原文件查看。",
    }


@app.get("/library/assets/{meeting_id}/{asset_id}")
def download_library_asset(meeting_id: str, asset_id: str, download: bool = False):
    asset, path = _find_archived_asset(meeting_id, asset_id)
    disposition = "attachment" if download else "inline"
    return FileResponse(path, filename=_safe_filename(asset.get("name")), content_disposition_type=disposition)


@app.get("/library/search")
def library_search(q: str = ""):
    query = q.strip().lower()
    if not query:
        return {"meetings": [], "actions": []}
    wildcard = f"%{query}%"
    with _db() as conn:
        meetings = [dict(row) for row in conn.execute("SELECT id, meeting_time, meeting_type, discipline, summary FROM meetings WHERE lower(summary) LIKE ? OR lower(data_json) LIKE ? ORDER BY meeting_time DESC LIMIT 30", (wildcard, wildcard))]
        actions = [dict(row) for row in conn.execute("SELECT * FROM action_items WHERE lower(description) LIKE ? OR lower(owner) LIKE ? ORDER BY updated_at DESC LIMIT 50", (wildcard, wildcard))]
    return {"meetings": meetings, "actions": actions}


class ActionStatusUpdate(BaseModel):
    status: Optional[Literal["open", "in_progress", "done", "blocked"]] = None
    response_status: Optional[Literal["unanswered", "responded"]] = None
    response_note: Optional[str] = None
    description: Optional[str] = None
    owner: Optional[str] = None
    due: Optional[str] = None
    level: Optional[str] = None
    risk: Optional[str] = None


class ActionItemCreate(BaseModel):
    meeting_id: str
    description: str
    reporter: str = "待确认"
    owner: str = "待确认"
    due: str = "未明确"
    level: str = "YELLOW"
    risk: str = ""


@app.post("/library/action-items")
def create_action_item(item: ActionItemCreate):
    _meeting_row_with_data(item.meeting_id)
    description = item.description.strip()
    if not description:
        raise HTTPException(status_code=400, detail="行动项内容不能为空")
    action_id, now = str(uuid.uuid4()), datetime.now().isoformat(timespec="seconds")
    with _db() as conn:
        conn.execute(
            """INSERT INTO action_items
            (id, meeting_id, reporter, description, owner, due, level, risk, status, evidence_json, response_status, response_note, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, 'open', '[]', 'unanswered', '', ?, ?)""",
            (action_id, item.meeting_id, item.reporter.strip() or "待确认", description, item.owner.strip() or "待确认", item.due.strip() or "未明确", item.level.strip() or "YELLOW", item.risk.strip(), now, now),
        )
    return {"id": action_id}


@app.patch("/library/action-items/{action_id}")
def update_action_status(action_id: str, update: ActionStatusUpdate):
    changes, values = [], []
    if update.status is not None:
        changes.append("status=?")
        values.append(update.status)
    if update.response_status is not None:
        changes.append("response_status=?")
        values.append(update.response_status)
    if update.response_note is not None:
        changes.append("response_note=?")
        values.append(update.response_note.strip()[:500])
    for field in ("description", "owner", "due", "level", "risk"):
        value = getattr(update, field)
        if value is not None:
            changes.append(f"{field}=?")
            values.append(value.strip()[:2000])
    if not changes:
        raise HTTPException(status_code=400, detail="请提供需要更新的行动项状态")
    now = datetime.now().isoformat(timespec="seconds")
    with _db() as conn:
        result = conn.execute(f"UPDATE action_items SET {', '.join(changes)}, updated_at=? WHERE id=?", (*values, now, action_id))
        if result.rowcount == 0:
            raise HTTPException(status_code=404, detail="未找到行动项")
    return {"id": action_id, "status": update.status, "response_status": update.response_status}


@app.delete("/library/action-items/{action_id}")
def delete_action_item(action_id: str):
    with _db() as conn:
        result = conn.execute("DELETE FROM action_items WHERE id=?", (action_id,))
        if result.rowcount == 0:
            raise HTTPException(status_code=404, detail="未找到行动项")
    return {"id": action_id, "deleted": True}


@app.get("/trial-case")
def trial_case_manifest():
    """Metadata for the one-click browser trial. Never expose local source paths."""
    return TRIAL_CASE


@app.get("/trial-case/assets/{asset_id}")
def download_trial_asset(asset_id: str):
    allowed = {item["id"]: item for item in TRIAL_CASE["materials"]}
    allowed[TRIAL_CASE["reference_download"]["id"]] = TRIAL_CASE["reference_download"]
    asset = allowed.get(asset_id)
    if not asset:
        raise HTTPException(status_code=404, detail="未找到试用材料")
    path = TRIAL_ASSETS_ROOT / asset["filename"]
    if not path.is_file():
        raise HTTPException(status_code=404, detail="试用材料尚未随应用部署")
    return FileResponse(path, filename=asset["filename"], content_disposition_type="attachment")


async def _process_meeting(request: Optional[Request] = None, progress=None, form=None):
    async def emit(stage: str, percent: int, message: str):
        if progress is not None:
            await progress({"event": "progress", "stage": stage, "percent": percent, "message": message})

    # 流式接口必须先在 HTTP 响应开始前读完 multipart；否则浏览器上传流会被截断。
    if form is None:
        if request is None:
            raise HTTPException(status_code=400, detail="缺少组会素材")
        form = await request.form()
    time = form.get("time") or ""
    if not time:
        raise HTTPException(status_code=400, detail="缺少 time")
    topic = form.get("topic") or ""
    roles = form.get("roles") or "[]"
    transcript = form.get("transcript") or ""
    transcript_original_input = form.get("transcript_original") or ""
    meeting_type = (form.get("meeting_type") or "").strip()
    discipline = (form.get("discipline") or "").strip()
    files_list = form.getlist("files") or []
    files_list = [x for x in files_list if getattr(x, "filename", None)]
    transcript_file = form.get("transcript_file")
    await emit("received", 8, "已接收素材，正在建立本次会议的处理任务…")
    def _parse_name_list(s: str) -> list:
        if not s or not s.strip():
            return []
        try:
            out = json.loads(s)
            if isinstance(out, list):
                return [str(x).strip() for x in out if str(x).strip()]
        except Exception:
            pass
        return [r.strip() for r in re.split(r"[,，、\s]+", s) if r.strip()]

    roles_list = _parse_name_list(roles)
    reporters_list = _parse_name_list(form.get("reporters") or "")
    speaker_materials = _parse_json_list(form.get("speaker_materials") or "[]")

    ppt_texts = []
    ppt_slides = []
    pdf_texts = []
    papers_pages = []  # PDF 按页：{ page, text, image }，page 为全局 1-based
    transcript_extra_parts = []
    uploaded_file_names = []  # 原始文件名，供模型从文件名推断汇报人（如「张三汇报.pptx」）
    source_documents = []  # 保留文件级身份，供证据链精确定位与前端预览
    archive_inputs = []  # 原始文件按「会议—人员」写入本地资料库

    def _material_owners(filename: str):
        owners = []
        for profile in speaker_materials:
            if not isinstance(profile, dict):
                continue
            if filename in (profile.get("materials") or []):
                name = str(profile.get("name") or "").strip()
                if name:
                    owners.append(name)
        return owners or ["其他材料"]

    file_list = list(files_list) if files_list else []
    if getattr(transcript_file, "filename", None):
        transcript_original_name = transcript_file.filename
        transcript_name = transcript_original_name.lower()
        transcript_bytes = await transcript_file.read()
        if transcript_name.endswith(".txt"):
            transcript_extra_parts.append(transcript_bytes.decode("utf-8", errors="replace"))
        elif transcript_name.endswith(".docx"):
            transcript_extra_parts.append(extract_docx_text_bytes(transcript_bytes))
        else:
            raise HTTPException(status_code=400, detail="转写稿仅支持 TXT 或 DOCX 格式。")
        archive_inputs.append({
            "name": transcript_original_name,
            "type": "transcript",
            "owners": ["会议转写"],
            "content": transcript_bytes,
            "preview_text": transcript_extra_parts[-1] if transcript_extra_parts else "",
        })
    for f in file_list:
        original_filename = (getattr(f, "filename") or "").strip()
        source_id = f"source_{len(source_documents) + 1}"
        uploaded_file_names.append(original_filename)
        name = original_filename.lower()
        try:
            if asyncio.iscoroutinefunction(getattr(f, "read", None)):
                content = await f.read()
            else:
                content = getattr(f, "file", None)
                content = content.read() if content else b""
        except Exception:
            content = b""
        if not isinstance(content, bytes):
            content = b""
        if name.endswith((".ppt", ".pptx")):
            archive_type = "ppt"
        elif name.endswith(".pdf"):
            archive_type = "pdf"
        elif name.endswith((".doc", ".docx")):
            archive_type = "word"
        else:
            archive_type = "file"
        archive_entry = {
            "name": original_filename,
            "type": archive_type,
            "owners": _material_owners(original_filename),
            "content": content,
            "preview_text": "",
        }
        if name.endswith((".ppt", ".pptx")):
            text, slides_meta = extract_ppt_content_bytes(content)
            archive_entry["preview_text"] = text
            ppt_texts.append(f"[文件 {source_id}：{original_filename}]\n{text}")
            page_images = render_ppt_to_page_images(content)
            if page_images:
                rendered_meta = []
                for idx, img in enumerate(page_images):
                    meta = dict(slides_meta[idx]) if idx < len(slides_meta) else {"page": idx + 1}
                    meta.update({"page": idx + 1, "images": [img], "source_id": source_id, "source_name": original_filename})
                    rendered_meta.append(meta)
                slides_meta = rendered_meta
            else:
                for slide in slides_meta:
                    slide["source_id"] = source_id
                    slide["source_name"] = original_filename
            ppt_slides.extend(slides_meta)
            source_documents.append({"id": source_id, "name": original_filename, "type": "ppt", "text": text})
        elif name.endswith(".pdf"):
            pdf_text = extract_pdf_text_bytes(content)
            archive_entry["preview_text"] = pdf_text
            pdf_texts.append(f"[文件 {source_id}：{original_filename}]\n{pdf_text}")
            pdf_pages = extract_pdf_pages_with_images(content)
            for p in pdf_pages:
                p["source_page"] = p.get("page")
                p["page"] = len(papers_pages) + 1
                p["source_id"] = source_id
                p["source_name"] = original_filename
                papers_pages.append(p)
            source_documents.append({"id": source_id, "name": original_filename, "type": "pdf", "text": pdf_text})
        elif name.endswith((".doc", ".docx")):
            docx_text = extract_docx_text_bytes(content)
            archive_entry["preview_text"] = docx_text
            if any(kw in name for kw in ("转写", "语音", "转文字", "会议记录", "纪要", "实录", "记录")):
                transcript_extra_parts.append(docx_text)
            else:
                pdf_texts.append(f"[文件 {source_id}：{original_filename}]\n{docx_text}")
                source_documents.append({"id": source_id, "name": original_filename, "type": "word", "text": docx_text})
        else:
            pass
        archive_inputs.append(archive_entry)

    await emit("parsed", 35, f"已解析 {len(source_documents)} 份材料，正在建立文件级证据索引…")

    transcript_full = transcript.strip()
    if transcript_extra_parts:
        transcript_full = (transcript_full + "\n\n" + "\n\n".join(transcript_extra_parts)).strip()

    transcript_original = transcript_original_input.strip() or transcript_full
    if transcript_original:
        archive_inputs.append({
            "name": "组会原始转写稿.txt",
            "type": "transcript",
            "owners": ["会议转写"],
            "content": transcript_original.encode("utf-8"),
            "preview_text": transcript_original,
        })

    ppt_combined = "\n\n".join(ppt_texts) if ppt_texts else ""
    pdf_combined = "\n\n".join(pdf_texts) if pdf_texts else ""

    # ── System Prompt ──────────────────────────────────────────────────────
    # 修改提示词请直接编辑下方。会议类型仅作参考，要点以本场会议实际讨论为准。
    # ──────────────────────────────────────────────────────────────────────

    MEETING_TYPE_GUIDES = {
        "人文社科": (
            "  1. \"核心论点\"：本次汇报的主要学术观点与论证脉络\n"
            "  2. \"论文结构与逻辑问题\"：写作结构、逻辑漏洞或老师指出的问题\n"
            "  3. \"文献综述要点\"：重要文献的讨论、引用建议与研究空白\n"
            "  4. \"修改意见与写作计划\"：明确的修改方向与下一步写作任务\n"
        ),
        "实验工科": (
            "  1. \"实验目的与方法\"：实验设计、研究问题与方法选择\n"
            "  2. \"数据与图表结果\"：关键实验数据、图表解读与主要结论\n"
            "  3. \"异常问题与原因分析\"：失败、异常结果及原因分析\n"
            "  4. \"下周实验安排\"：后续实验计划与参数调整\n"
        ),
        "理论数理": (
            "  1. \"研究思路与问题定义\"：核心研究问题与整体思路\n"
            "  2. \"公式/模型/算法要点\"：关键数学推导、模型结构或算法设计\n"
            "  3. \"当前难点\"：尚未解决的理论障碍或推导瓶颈\n"
            "  4. \"下一步推导方向\"：老师建议的后续理论探索方向\n"
        ),
        "通用组会": (
            "  1. \"本周完成情况\"：本周已完成的任务与进展\n"
            "  2. \"遇到的问题\"：本周遇到的困难、阻碍与未解决事项\n"
            "  3. \"老师意见与反馈\"：导师在会上给出的具体评价与建议\n"
            "  4. \"下周计划\"：下周的明确任务目标与时间节点\n"
        ),
    }

    type_guide_block = ""
    for t, guide in MEETING_TYPE_GUIDES.items():
        type_guide_block += f'\n\u25b7 meeting_type = "{t}" \u65f6\uff0c\u53ef\u4f9d\u6b21\u4f5c\u4e3a\u5173\u6ce8\u70b9\uff08\u4ec5\u53c2\u8003\uff09\uff1a\n{guide}'

    user_meeting_type = (meeting_type or "").strip()
    user_discipline = (discipline or "").strip()

    role_instruction = ""
    if user_meeting_type and user_meeting_type in MEETING_TYPE_GUIDES:
        role_instruction = (
            "\n【重要：以用户选择为准】用户已选择会议类型为「" + user_meeting_type + "」。\n"
            "- meeting_type 必须且只能使用该值，不要根据转写/PPT/文献内容自行推断或改写。\n"
            "- 请以该类型领域专家的身份撰写纪要，使用该领域常用术语与规范，使输出更具针对性和专业性。\n\n"
        )
    if user_discipline:
        role_instruction = (
            role_instruction
            + "\n【重要：以用户填写学科为准】用户指定了学科「" + user_discipline + "」。\n"
            "- 请以该学科领域专家身份撰写纪要，使用该学科常用术语、关注该领域典型问题与规范，使输出更具针对性和专业性。\n"
            "- 不要根据内容猜测学科或类型，全文须体现该学科视角。若用户未选择会议类型，meeting_type 可标为「通用组会」，但内容须紧扣该学科。\n\n"
        )
    if reporters_list:
        role_instruction = (
            role_instruction
            + "\n【重要：按用户指定汇报人展开，不得遗漏、不得删除】用户指定了本次汇报人名单：" + "、".join(reporters_list) + "。\n"
            "- by_reporter 必须严格按该名单顺序与人数输出：名单里有几人就输出几条，顺序与名单一致。\n"
            "- 二、核心讨论要点 与 三、导师反馈与行动项 均按每位汇报人分别展开；每条对应一名汇报人的 reporter、key_points、advisor_feedback、action_items。\n"
            "- 若某汇报人在素材中未识别到其汇报内容，该汇报人仍须保留一条，其 key_points、advisor_feedback、action_items 可为空数组，或写一条 key_points：title 为「未识别到该汇报人相关内容」、detail 为「（根据当前素材未匹配到该汇报人的发言或材料）」、evidence 为空数组。不得直接删掉该汇报人条目。\n"
            "- utterances 中 speaker 若能从转写中对应到上述姓名，请尽量使用用户提供的汇报人姓名，便于与 by_reporter 一致。\n\n"
        )
    if speaker_materials:
        role_instruction = (
            role_instruction
            + "\n【汇报人与材料绑定】用户已显式标注每位发言人关联的材料。该绑定是高优先级上下文："
            + json.dumps(speaker_materials, ensure_ascii=False)
            + "。\n- 优先使用该汇报人绑定材料与其发言交叉印证；不得把另一位汇报人的材料归到其名下。\n"
            + "- 材料绑定只能辅助归属，不能替代原始证据；无明确内容时须明确标注信息不足。\n\n"
        )

    system_prompt = (
        "【研行记 · 身份与定位】\n"
        "你是「研行记 (Research Action Note)」的核心引擎，面向科研组会场景，负责把语音转写、PPT、文献等多源素材加工成「高价值、可执行、可溯源」的组会纪要。\n"
        "用途：帮助师生在会后快速回顾讨论要点、导师反馈与行动项，并能在纪要中追溯到「谁在何时说了什么、对应哪一页材料」，从而减少信息遗漏、明确下一步任务。\n"
        "特点：多来源证据链（每条要点/反馈/行动项尽量挂接 transcript、PPT 页、文献片段）；按汇报人分块组织；导师说的话要细致保留、不合并不省略；会议类型仅作参考，具体表述以本场会议实际讨论内容为准。\n"
        "语气与风格：专业、克制、信息密度高。用学术化、正式的中文表述，不夸张、不抒情；重点放在「事实与建议」的准确传达，让读者读完即知本场组会讨论了什么、老师具体提了哪些意见、接下来要做什么。\n\n"

        "【专业质量底线】\n"
        "- 严格区分三类信息：已发生的事实/数据、发言者或导师的判断、模型基于材料的整理。没有证据时不补写实验结果、文献结论、时间节点或因果关系。\n"
        "- 仅在导师或发言者明确提出时，才把内容写为「要求」「结论」或 RED 行动项；学生的设想、待验证假设要明确写为「建议」或「待验证」。\n"
        "- 数值、样品/参数、方法、作者与文献名必须逐字核对来源；无法核对时宁可省略，不得似是而非地改写。\n"
        "- 每条行动项必须尽量给出 owner（负责人）、due（明确截止时间；未说明则填「未明确」）和 risk（完成风险/依赖；无则填「无明显风险」）。\n\n"

        "【输出形式】\n"
        "请严格按以下阶段工作，最终只输出一个 JSON 对象，不要输出中间思考过程。\n\n"
        + role_instruction +

        "【阶段0：会议类型】\n"
        "仅当用户未选择类型且未填写学科时，才根据转写、PPT、文献判断 meeting_type；否则必须以上方用户选择/填写的为准。\n"
        "类型说明（供参考）：\"人文社科\"多涉及论文写作、文献综述、理论框架；\"实验工科\"多涉及实验设计、数据与图表、工程排查；\"理论数理\"多涉及数学推导、公式建模、算法设计；\"通用组会\"为周期性进展汇报或难以归类的讨论。\n"
        "注意：会议类型只是标签与参考，不约束你如何组织要点。本场会议究竟讨论了什么，就写什么；要点条数、小标题可根据实际内容灵活增减，不必拘泥于类型下的固定几条。\n\n"

        "【阶段1：汇报人多源识别与语音角色梳理】\n"
        "汇报人可从以下多源综合识别，请结合使用：\n"
        "1) 用户填写的「汇报人」名单（若提供则优先，且 by_reporter 须严格按该名单展开、不删不漏）；\n"
        "2) 上传文件名：meeting_info 中的 uploaded_file_names 可能含汇报人姓名（如「张三汇报.pptx」「李四组会.pdf」），可从中推断谁汇报；\n"
        "3) PPT 正文：ppt_text 开头常为标题页或汇报人信息，可结合识别；\n"
        "4) 语音转写：将 transcript 拆分为 utterances，每条含 index、speaker、content，speaker 尽量区分「导师/老师」与汇报人；转写中出现的发言人姓名、称谓可作为汇报人依据。\n"
        "若出现多位学生/汇报人，不得合并为一人；speaker 命名需稳定，尽量与用户提供的汇报人姓名或文件名中的姓名一致。导师发言不得混入汇报人发言中。\n"
        "说明：系统会用 utterances 做纪要生成；溯源时仍展示用户提供的原始转写原文，便于对照当时原话。\n\n"

        "【阶段2：素材对齐与证据链构建（多来源）】\n"
        "每条 key_points、advisor_feedback、action_items 至少一条 evidence，尽量多来源（转写、PPT 页、文献片段）。\n"
        "evidence 含 type(transcript|ppt|paper)、source_id、source_name、page、location、quote、context、note。对 PPT/PDF/Word 必须从 source_documents 中原样拷贝 source_id 与 source_name；page 为页码（Word 可省略）；location 也要写清具体出处。不得捏造页码、文件名或原文。\n"
        "本次消息可能额外附带 PPT/PDF 的页面图像；它们和 visual_materials 中的 source_id、页码一一对应。图表、公式、版式或扫描页中读到的信息必须用对应 ppt/paper evidence 标注，不能把视觉材料误写为转写证据。\n"
        "重要：对 type=transcript 的 evidence，quote 必须为转写原文中的连续片段（可模糊对应），以便在溯源界面高亮显示并上下滚动查看上下文。\n\n"

        "【阶段3：按汇报人组织 + 要点结构】\n"
        "输出 by_reporter，每人含 reporter、key_points、advisor_feedback、action_items。\n"
        "若用户提供了汇报人名单：by_reporter 必须按该名单顺序、每人一条；未识别到内容的汇报人保留条目，key_points/advisor_feedback/action_items 留空或写「未识别到该汇报人相关内容」，不得删除。\n"
        "若用户未提供名单：从 transcript、文件名、PPT 等识别汇报人，by_reporter 覆盖所有识别到的汇报人，不得只输出部分。\n"
        "「本次组会要点」对应 key_points。下面按会议类型给了一些可参考的方面，仅供你组织时参考，不必严格照搬：可根据本场会议实际讨论了什么来定条数与标题。重点是把「本场会议真实讨论到的内容」写全、写细。\n"
        "每位有内容的汇报人至少输出 1 条 key_points；若某汇报人信息不足且用户未要求保留，可写“该汇报人本次有效信息较少”；若用户名单中有此人则必须保留其条目且不得删除。\n"
        + type_guide_block + "\n\n"

        "【导师/老师说的内容：务必细致、不遗漏】\n"
        "导师（或老师）在会上的发言是纪要的核心价值之一，必须细致呈现，不要漏点、不要合并成一句笼统话。\n"
        "- 每条独立的意见、建议、批评、肯定，尽量单独成条或单独成句写清，避免「老师提了几点意见」这种概括。\n"
        "- 老师提到的具体信息必须保留：文献名、作者、年份、书名、论文题目、方法名、工具名、数据来源、某章某节、某页某段等，在 advisor_feedback 的 content 或 action_items 的 description 中写清楚。例如：「导师建议阅读张三等(2020)《某某研究》中关于……的论述」「可参考 XX 方法/XX 数据集」「建议先做 A 再做 B」。\n"
        "- 老师指出的具体问题（如某处逻辑、某段表述、某组实验）要写清是「哪方面的问题」以及「老师的大致意见」，不要只写「老师认为需要修改」而丢失具体指向。\n"
        "- 老师给出的下一步任务、时间节点、交付物，要在 action_items 中写具体：做什么、做到什么程度、为什么（reason 可简要说明）。\n"
        "若转写或 PPT 中已有上述具体表述，请原样或适度整理后写入对应条目的 content/description，确保读者仅凭纪要就能还原老师的主要意见与要求。\n\n"

        "【行动项等级与表述】\n"
        "action_items 的 level：RED=必须完成、有明确期限或硬性要求；YELLOW=建议优化、应优先考虑；GREEN=探索性、可选尝试。description 写清「做什么」；reason 可简要写「为何被判定为该等级或依据哪句讨论」。\n\n"

        "【导师关注点】\n"
        "输出 mentor_focus：3–8 个导师本场反复强调或明确追问的短语，如「对照实验」「误差条」「机制证据」。每个 advisor_feedback 可选 focus_tags（0–3 个），只能使用会议中真实出现的关注点，不要编造。\n\n"

        "【专属视角与跨领域启发】\n"
        "补充以下字段，全部严格依据本场材料，信息不足时返回空数组或空字符串：\n"
        "- mentor_brief：120 字内的导师摘要，提炼导师的判断框架、最重要的风险和下一步优先级；不是泛泛复述。\n"
        "- student_progress：逐位汇报人输出 student、progress（本场相对推进/卡点）、next_focus（下次最应回应的一个点）。\n"
        "- term_glossary：3–8 个材料中实际出现且对跨专业读者不直观的术语，每项含 term、explanation（50字内、通俗解释）、evidence。\n"
        "- innovation_transfers：0–3 条跨领域迁移启发，每项含 source_domain、target_domain、transferable_idea、application、caution、evidence。只能在材料确有可迁移方法/机制时提出，明确标注为「启发/待验证」，不能写成已证实结论。\n"
        "- socratic_questions：3 条下一次组会可用的追问，每项含 question、why（它要检验的假设/证据缺口）、target（建议回答人）。问题必须具体、可回答，不能是空泛的「请深入思考」。\n\n"

        "【JSON 输出结构（必须含 by_reporter）】\n"
        "{\n"
        "  \"meeting_type\": \"人文社科 | 实验工科 | 理论数理 | 通用组会\",\n"
        "  \"basic_info\": { \"time\": \"...\", \"topic\": \"...\", \"roles\": [\"...\"] },\n"
        "  \"utterances\": [ { \"index\": 1, \"speaker\": \"导师\", \"content\": \"...\" } ],\n"
        "  \"by_reporter\": [\n"
        "    {\n"
        "      \"reporter\": \"汇报人姓名\",\n"
        "      \"key_points\": [ { \"title\": \"...\", \"detail\": \"...\", \"evidence\": [ { \"type\": \"...\", \"source_id\": \"source_1\", \"source_name\": \"...\", \"page\": 3, \"location\": \"...\", \"quote\": \"...\", \"context\": \"...\", \"note\": \"...\" } ] } ],\n"
        "      \"advisor_feedback\": [ { \"speaker\": \"导师\", \"content\": \"...\", \"focus_tags\": [\"...\"], \"evidence\": [ ... ] } ],\n"
        "      \"action_items\": [ { \"level\": \"RED|YELLOW|GREEN\", \"description\": \"...\", \"owner\": \"...|待确认\", \"due\": \"...|未明确\", \"risk\": \"...\", \"reason\": \"...\", \"evidence\": [ ... ] } ]\n"
        "    }\n"
        "  ],\n"
        "  \"summary\": \"150字内总结\",\n"
        "  \"mentor_focus\": [\"导师高频关注点\"],\n"
        "  \"mentor_brief\": \"导师摘要\",\n"
        "  \"student_progress\": [ { \"student\": \"汇报人\", \"progress\": \"...\", \"next_focus\": \"...\" } ],\n"
        "  \"term_glossary\": [ { \"term\": \"术语\", \"explanation\": \"通俗解释\", \"evidence\": [ ... ] } ],\n"
        "  \"innovation_transfers\": [ { \"source_domain\": \"来源领域\", \"target_domain\": \"迁移领域\", \"transferable_idea\": \"可迁移机制/方法\", \"application\": \"本研究可尝试方式\", \"caution\": \"待验证边界\", \"evidence\": [ ... ] } ],\n"
        "  \"socratic_questions\": [ { \"question\": \"具体追问\", \"why\": \"要检验什么\", \"target\": \"建议回答人\" } ]\n"
        "}\n\n"

        "【阶段4：表达与格式】\n"
        "- 几乎全部使用中文，必要时可以使用英文术语，措辞学术、正式、信息密度高。basic_info 从 meeting_info 拷贝或从素材中提炼。\n"
        "- by_reporter 与 summary 严格依据 transcript、ppt_text、papers_text 生成；不编造未出现的讨论；每条要点/反馈/行动项的 evidence 尽量覆盖多类来源。\n"
        "- 导师反馈与行动项中凡有具体文献、方法、问题指向、时间节点、交付要求的，一律写清，不省略、不泛化为「按老师意见修改」等笼统表述。\n"
    )

    user_content = {
        "meeting_info": {
            "time": time,
            "topic": topic,
            "roles": roles_list,
            "reporters": reporters_list,
            "speaker_materials": speaker_materials,
            "uploaded_file_names": uploaded_file_names,
        },
        "user_meeting_type": user_meeting_type or None,
        "user_discipline": user_discipline or None,
        "transcript": transcript_full,
        "ppt_text": ppt_combined,
        "papers_text": pdf_combined,
        "source_documents": [{"id": item["id"], "name": item["name"], "type": item["type"]} for item in source_documents],
    }

    # 文字提取覆盖可复制文本；页面图片补足图表、公式、扫描件和 PPT 版式信息。
    visual_inputs, visual_materials = [], []
    if VISION_INPUT_ENABLED and VISION_MAX_IMAGES:
        for slide in ppt_slides:
            if len(visual_inputs) >= VISION_MAX_IMAGES:
                break
            images = slide.get("images") if isinstance(slide, dict) else []
            image = images[0] if isinstance(images, list) and images else None
            if not image:
                continue
            visual_inputs.append({"type": "image_url", "image_url": {"url": image}})
            visual_materials.append({
                "type": "ppt", "source_id": slide.get("source_id"), "source_name": slide.get("source_name"), "page": slide.get("page"),
            })
        for page in papers_pages:
            if len(visual_inputs) >= VISION_MAX_IMAGES:
                break
            image = page.get("image") if isinstance(page, dict) else None
            if not image:
                continue
            visual_inputs.append({"type": "image_url", "image_url": {"url": image}})
            visual_materials.append({
                "type": "paper", "source_id": page.get("source_id"), "source_name": page.get("source_name"), "page": page.get("page"),
            })
    user_content["visual_materials"] = visual_materials
    await emit("indexed", 52, f"证据索引已建立：文本材料 {len(source_documents)} 份，视觉页面 {len(visual_inputs)} 页。")

    user_message_content = [{"type": "text", "text": json.dumps(user_content, ensure_ascii=False)}]
    user_message_content.extend(visual_inputs)

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_message_content},
    ]

    try:
        await emit("model", 62, f"正在调用多模态模型 {MODEL_NAME}，分析文字、图表与页面内容…")
        response = await asyncio.to_thread(
            client.chat.completions.create,
            model=MODEL_NAME,
            messages=messages,
            response_format={"type": "json_object"},
        )
        content = response.choices[0].message.content
        data = json.loads(content)
        await emit("model_complete", 84, "模型已完成分析，正在校验证据链、行动项与专属洞察…")
    except Exception as exc:
        provider_status = getattr(exc, "status_code", None)
        if provider_status == 401:
            detail = "模型服务鉴权失败，请检查 OPENAI_API_KEY。"
        elif provider_status == 403:
            detail = f"模型“{MODEL_NAME}”当前不可用或不支持所在地区，请在 .env 中更换 OPENAI_MODEL 后重试。"
        elif provider_status == 429:
            detail = "模型服务请求过于频繁或额度不足，请稍后重试或检查账户额度。"
        else:
            detail = "模型生成失败，请检查网络、模型配置和服务状态后重试。"
        raise HTTPException(status_code=502, detail=detail) from exc

    data["model_used"] = MODEL_NAME
    if user_meeting_type and user_meeting_type in MEETING_TYPE_GUIDES:
        data["meeting_type"] = user_meeting_type

    data["raw_sources"] = {
        "transcript": transcript_full,
        "transcript_original": transcript_original,
        "ppt_text": ppt_combined,
        "papers_text": pdf_combined,
        "ppt_slides": ppt_slides,
        "papers_pages": papers_pages,
        "source_documents": source_documents,
    }

    def ensure_evidence(obj, fallback_quote: str):
        if not isinstance(obj, dict):
            return
        ev = obj.get("evidence")
        if not ev or not isinstance(ev, list) or len(ev) == 0:
            snippet = (transcript_full or "")[:400].strip() or "(无转写)"
            obj["evidence"] = [{
                "type": "transcript",
                "location": "语音转写",
                "quote": (fallback_quote or obj.get("content") or obj.get("detail") or obj.get("description") or "")[:200] or "见上文",
                "context": snippet,
                "note": "根据转写整理",
            }]

    for block in data.get("by_reporter") or []:
        for kp in block.get("key_points") or []:
            ensure_evidence(kp, kp.get("detail") or kp.get("title"))
        for fb in block.get("advisor_feedback") or []:
            ensure_evidence(fb, fb.get("content"))
        for ai in block.get("action_items") or []:
            ensure_evidence(ai, ai.get("description"))

    for kp in data.get("key_points") or []:
        ensure_evidence(kp, kp.get("detail") or kp.get("title"))
    for fb in data.get("advisor_feedback") or []:
        ensure_evidence(fb, fb.get("content"))
    for ai in data.get("action_items") or []:
        ensure_evidence(ai, ai.get("description"))

    data["mentor_focus"] = [str(term).strip() for term in (data.get("mentor_focus") or []) if str(term).strip()][:8]
    data["mentor_brief"] = str(data.get("mentor_brief") or "").strip()[:500]

    def _normalise_dict_list(value, fields, limit):
        if not isinstance(value, list):
            return []
        out = []
        for item in value:
            if not isinstance(item, dict):
                continue
            cleaned = {field: str(item.get(field) or "").strip()[:500] for field in fields}
            if any(cleaned.values()):
                if isinstance(item.get("evidence"), list):
                    cleaned["evidence"] = item["evidence"][:3]
                out.append(cleaned)
            if len(out) >= limit:
                break
        return out

    data["student_progress"] = _normalise_dict_list(data.get("student_progress"), ["student", "progress", "next_focus"], 12)
    data["term_glossary"] = _normalise_dict_list(data.get("term_glossary"), ["term", "explanation"], 8)
    data["innovation_transfers"] = _normalise_dict_list(data.get("innovation_transfers"), ["source_domain", "target_domain", "transferable_idea", "application", "caution"], 3)
    data["socratic_questions"] = _normalise_dict_list(data.get("socratic_questions"), ["question", "why", "target"], 5)
    for item in data["term_glossary"] + data["innovation_transfers"]:
        ensure_evidence(item, item.get("explanation") or item.get("transferable_idea"))

    def _is_teacher_like(name: str) -> bool:
        if not name:
            return False
        n = str(name).strip().lower()
        teacher_keys = ["导师", "老师", "教授", "pi", "老板", "导师a", "导师b", "teacher"]
        return any(k in n for k in teacher_keys)

    def _norm_name(name: str) -> str:
        return "".join(str(name or "").strip().lower().split())

    def _block_matches_reporter(blk: dict, reporter_name: str) -> bool:
        if not blk or not isinstance(blk, dict):
            return False
        r = str(blk.get("reporter") or "").strip()
        if not r:
            return False
        return _norm_name(r) == _norm_name(reporter_name)

    # 若用户提供了汇报人名单：严格按名单顺序保证每人一条，不删不漏；未识别到内容的留空块
    by_blocks = data.get("by_reporter")
    if not isinstance(by_blocks, list):
        by_blocks = []

    if reporters_list:
        ordered_blocks = []
        used_indices = set()
        for user_reporter in reporters_list:
            user_reporter = str(user_reporter or "").strip()
            if not user_reporter:
                continue
            matched = None
            for i, blk in enumerate(by_blocks):
                if i in used_indices:
                    continue
                if _block_matches_reporter(blk, user_reporter):
                    matched = blk
                    used_indices.add(i)
                    break
            if matched is not None:
                ordered_blocks.append(matched)
            else:
                ordered_blocks.append({
                    "reporter": user_reporter,
                    "key_points": [{
                        "title": "未识别到该汇报人相关内容",
                        "detail": "（根据当前素材未匹配到该汇报人的发言或材料）",
                        "evidence": [],
                    }],
                    "advisor_feedback": [],
                    "action_items": [],
                })
        if len(ordered_blocks) < len([r for r in reporters_list if str(r).strip()]):
            seen = {_norm_name(str(b.get("reporter") or "")) for b in ordered_blocks}
            for user_reporter in reporters_list:
                user_reporter = str(user_reporter or "").strip()
                if not user_reporter or _norm_name(user_reporter) in seen:
                    continue
                ordered_blocks.append({
                    "reporter": user_reporter,
                    "key_points": [{
                        "title": "未识别到该汇报人相关内容",
                        "detail": "（根据当前素材未匹配到该汇报人的发言或材料）",
                        "evidence": [],
                    }],
                    "advisor_feedback": [],
                    "action_items": [],
                })
                seen.add(_norm_name(user_reporter))
        by_blocks = ordered_blocks
    else:
        existing_reporters = set()
        for blk in by_blocks:
            if not isinstance(blk, dict):
                continue
            r = str(blk.get("reporter") or "").strip()
            if r and (not _is_teacher_like(r)):
                existing_reporters.add(_norm_name(r))

        utterances = data.get("utterances") or []
        speaker_stats = {}
        for u in utterances:
            if not isinstance(u, dict):
                continue
            spk = str(u.get("speaker") or "").strip()
            content_u = str(u.get("content") or "").strip()
            if not spk or not content_u:
                continue
            if _is_teacher_like(spk):
                continue
            key = _norm_name(spk)
            if not key:
                continue
            if key not in speaker_stats:
                speaker_stats[key] = {"name": spk, "turns": 0, "chars": 0, "samples": []}
            speaker_stats[key]["turns"] += 1
            speaker_stats[key]["chars"] += len(content_u)
            if len(speaker_stats[key]["samples"]) < 2:
                speaker_stats[key]["samples"].append(content_u[:120])

        missing = []
        for key, st in speaker_stats.items():
            if (st["turns"] >= 2 or st["chars"] >= 60) and key not in existing_reporters:
                missing.append(st)

        for st in missing:
            sample_text = "；".join(st["samples"]).strip() or "该汇报人本次有效信息较少。"
            by_blocks.append({
                "reporter": st["name"],
                "key_points": [{
                    "title": "发言要点（自动补全）",
                    "detail": sample_text,
                    "evidence": [{
                        "type": "transcript",
                        "location": "语音转写",
                        "quote": sample_text[:120],
                        "context": (transcript_full or "")[:400] or sample_text,
                        "note": "根据转写自动补齐遗漏汇报人",
                    }]
                }],
                "advisor_feedback": [],
                "action_items": [],
            })

    if by_blocks:
        data["by_reporter"] = by_blocks

    chosen_type = (meeting_type or "").strip()
    chosen_discipline = (discipline or "").strip()
    if chosen_type in MEETING_TYPE_GUIDES:
        data["meeting_type"] = chosen_type
        data["user_selected_meeting_type"] = chosen_type
    if chosen_discipline:
        data["user_selected_discipline"] = chosen_discipline
        if chosen_type not in MEETING_TYPE_GUIDES:
            data["meeting_type"] = data.get("meeting_type") or "通用组会"

    await emit("saving", 94, "正在归档原始材料，并写入跨组会行动项与关注点…")
    data["library_meeting_id"] = _save_meeting(data, time, chosen_discipline, archive_inputs)
    await emit("complete", 100, "纪要、证据链与资料库已生成完成。")
    return data


@app.post("/process-meeting")
async def process_meeting(request: Request):
    """兼容原有前端的普通 JSON 接口。"""
    return await _process_meeting(request)


@app.post("/process-meeting-stream")
async def process_meeting_stream(request: Request):
    """将后端真实处理节点实时推送给前端；不以计时器伪造进度。"""
    try:
        # 先完整消费 multipart 请求体，再创建 SSE 响应，避免文件边界与响应流并发冲突。
        form = await request.form()
    except Exception as exc:
        raise HTTPException(status_code=400, detail="上传素材解析失败，请重新选择文件后重试。") from exc
    queue = asyncio.Queue()

    async def report(payload: dict):
        await queue.put(payload)

    async def run():
        try:
            data = await _process_meeting(progress=report, form=form)
            await queue.put({"event": "result", "data": data})
        except HTTPException as exc:
            await queue.put({"event": "error", "detail": exc.detail, "status": exc.status_code})
        except Exception:
            logger.exception("Meeting stream failed")
            await queue.put({"event": "error", "detail": "处理失败，请检查模型配置与素材格式后重试。", "status": 500})

    async def event_stream():
        task = asyncio.create_task(run())
        try:
            while True:
                payload = await queue.get()
                yield f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"
                if payload.get("event") in {"result", "error"}:
                    break
        finally:
            if not task.done():
                task.cancel()
                await asyncio.gather(task, return_exceptions=True)

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )
